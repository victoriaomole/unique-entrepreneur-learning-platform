from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def get_project_details():
    print("=== Kickoff Deck Generator ===")
    client_name = input("Client / Organisation Name: ").strip() or "Client"
    project_name = input("Project Name: ").strip() or "Project"
    phase_name = input("Phase Name (e.g. 'Phase 1 Kick-off'): ").strip() or "Phase 1 Kick-off"
    ba_name = input("Your Name (BA / Lead): ").strip() or "Omole Victoria Oluwatosin"

    # basic targets (you can tune per run)
    mvp_month = input("MVP Target (e.g. 'Month 7') [press Enter to use 'Month 7']: ").strip() or "Month 7"
    pilot_orgs = input("Pilot orgs target [default 3]: ").strip() or "3"
    completion_rate = input("Completion rate target % [default 80]: ").strip() or "80"
    paid_courses = input("Paid offerings / courses target [default 5]: ").strip() or "5"

    # build output filename
    safe_client = client_name.replace(" ", "")
    safe_phase = phase_name.replace(" ", "")
    output_filename = f"{safe_phase}_Kickoff_{safe_client}.pptx"

    return {
        "client": client_name,
        "project": project_name,
        "phase": phase_name,
        "ba": ba_name,
        "mvp_month": mvp_month,
        "pilot_orgs": pilot_orgs,
        "completion_rate": completion_rate,
        "paid_courses": paid_courses,
        "output_filename": output_filename,
    }


def build_slides(cfg):
    return [
        (
            f"{cfg['phase']} – {cfg['project']}",
            "Agenda:\n"
            "1. Vision & Background\n"
            "2. Goals & Success Criteria\n"
            "3. Scope Overview\n"
            "4. Roles & Responsibilities\n"
            "5. Discovery Plan & Deliverables\n"
            "6. Risks & Next Steps"
        ),
        (
            "Project Vision",
            f"To design and deliver a scalable, secure learning / data platform for {cfg['client']},\n"
            "enabling governed, insight-driven experiences for learners and stakeholders."
        ),
        (
            "Problem Statement / Opportunity",
            "Current landscape:\n"
            "- Fragmented tools and manual processes\n"
            "- Limited visibility across tenants / teams\n"
            "- Governance and compliance overhead\n\n"
            "Opportunity:\n"
            "- One governed platform that unifies delivery, reporting, and experience."
        ),
        (
            "Goals & Success Metrics",
            f"- MVP live by {cfg['mvp_month']} – ≥ {cfg['pilot_orgs']} pilot organisations\n"
            f"- Completion rate ≥ {cfg['completion_rate']} %\n"
            f"- ≥ {cfg['paid_courses']} monetised offerings (if applicable)\n"
            "- DPIA / GDPR / security checks passed pre-launch"
        ),
        (
            "Scope Overview (MVP)",
            "In Scope:\n"
            "- Tenant / organisation hierarchy & access control\n"
            "- Core feature set (content, enrolment, tracking)\n"
            "- Essential reports & dashboards\n"
            "- Secure authentication / roles\n\n"
            "Out of Scope (Phase 1):\n"
            "- Advanced AI and non-priority integrations\n"
            "- Full mobile/offline build"
        ),
        (
            "Roles & Responsibilities",
            f"- Business Analysis – {cfg['ba']}\n"
            f"- Data & Reporting – {cfg['ba']}\n"
            f"- Data Governance & Compliance – {cfg['ba']}\n"
            "- Product / Engineering – To be confirmed with {cfg['client']}"
        ),
        (
            "Discovery Plan (3 Weeks)",
            "Week 1 – Stakeholders, current-state, risks\n"
            "Week 2 – Interviews, journeys, requirements\n"
            "Week 3 – Target vision, MVP, roadmap, sign-off"
        ),
        (
            "Governance & Compliance Preview",
            "Focus Areas:\n"
            "- Ownership, roles & responsibilities\n"
            "- Access controls & segregation of data\n"
            "- DPIA / GDPR, retention, consent\n"
            "- Logging, audit, incident response"
        ),
        (
            "Analytics & KPI Preview",
            "Example KPIs:\n"
            "- Active users by tenant\n"
            "- Enrolment vs completion\n"
            "- Engagement by course / cohort\n"
            "- Support tickets & issues\n"
            "- Adoption trend over time"
        ),
        (
            "Initial Risks & Assumptions",
            "Examples (customise per client):\n"
            "- Access to existing systems / data delayed\n"
            "- Limited SME availability for discovery\n"
            "- Integration complexity higher than expected\n"
            "- Assumption: scope disciplined around MVP"
        ),
        (
            "Next Steps & Action Owners",
            f"- Circulate minutes & deck – {cfg['ba']} (+1 day)\n"
            "- Confirm stakeholders & interview slots – Client / BA (+2–3 days)\n"
            "- Share existing documentation & system access – Client (+3 days)\n"
            "- Agree MVP scope & timeline – Joint sign-off"
        ),
        (
            "Q & A / Wrap-Up",
            "Confirm shared understanding\n"
            "Capture decisions and open points\n"
            "Formally launch Discovery / Phase"
        ),
    ]


def generate_kickoff_pptx(cfg):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blue_grey = RGBColor(30, 55, 90)
    white = RGBColor(255, 255, 255)

    slides = build_slides(cfg)

    for title, body in slides:
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = white

        slide.shapes.title.text = title

        body_placeholder = slide.placeholders[1]
        body_placeholder.text = body

        # footer
        left = Inches(0.5)
        top = Inches(6.8)
        width = Inches(12)
        height = Inches(0.4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        p = textbox.text_frame.add_paragraph()
        p.text = (
            f"{cfg['ba']} — Business Analyst | Data / Reporting | Data Governance"
        )
        p.font.size = Pt(10)
        p.font.color.rgb = blue_grey

    prs.save(cfg["output_filename"])
    print(f"\n✅ PowerPoint file created: {cfg['output_filename']}")


if __name__ == "__main__":
    config = get_project_details()
    generate_kickoff_pptx(config)




